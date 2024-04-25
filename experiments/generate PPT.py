from pptx import Presentation
import collections
import collections.abc

prs = Presentation()

Layout = prs.slide_layouts[1]

first_slide = prs.slides.add_slide(Layout)  # Adding first slide

first_slide.shapes.title.text = "Created by Miraalfasa"

first_slide.placeholders[1].text = "Created powerpoint using Python for experiment purpose"
# image_placeholder = first_slide.placeholders[1]
# image_placeholder.placeholder_format.type = 'PICTURE'
# print(image_placeholder)
# image = image_placeholder.insert_picture(r'C:\Users\admin\Downloads\langchain-ask-pdf-main\experiments\img.png')
# image_placeholder.
# first_slide.placeholders[1].image = ""
prs.save("First_presentation.pptx")
