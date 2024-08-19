from openai import OpenAI
import streamlit as st
import os
import openai
from dotenv import load_dotenv

import re
from pexelsapi.pexels import Pexels
from pptx import Presentation
from pytube import YouTube

load_dotenv()
openai.api_key = os.environ['OPENAI_API_KEY']
api = os.environ['PEXELS_API_KEY']

st.set_page_config(
    page_title="Chatbot",
    page_icon="💬"
)

topic = ""


def parse_response_to_ppt_content(response_to_parse):
    text = response_to_parse

    # Splitting text into sections
    sections = re.split(r'_+\n?', text.strip())

    # List to store dictionaries
    result = []

    # Regular expressions to extract title, image, and details
    # file_name_pattern = re.compile(r'PPT File Name:\s*(.*?)\s*Title:', re.IGNORECASE | re.DOTALL)
    # print(file_name_pattern)
    title_pattern = re.compile(r'Title:\s*(.*?)\s*Image:', re.IGNORECASE | re.DOTALL)
    image_pattern = re.compile(r'Image:\s*\[(.*?)\]\s*Details:', re.IGNORECASE | re.DOTALL)
    details_pattern = re.compile(r'Details:\s*(.*?)$', re.IGNORECASE | re.DOTALL)

    # Extracting information for each section
    for section in sections:
        # Extracting title, image, and details
        # file_name_match = re.search(file_name_pattern, section)
        # print("__________file_name_match__________")
        # print(file_name_match)
        title_match = re.search(title_pattern, section)
        # print("__________title_match__________")
        # print(title_match)
        image_match = re.search(image_pattern, section)
        # print("_________________image_match________________")
        # print(image_match)
        details_match = re.search(details_pattern, section)
        # print("_______________details_match________________")
        # print(details_match)

        # Creating dictionary for the section
        if title_match and image_match and details_match:
            # print("___________________title_match.group(1).strip()_______________________")
            # print(title_match.group(1).strip())
            section_dict = {
                # "PPT File Name": file_name_match.group(1).strip(),
                "Title": title_match.group(1).strip(),
                "Image": image_match.group(1).strip(),
                "Details": details_match.group(1).strip()
            }
            result.append(section_dict)

    return result


def response_parser_response_to_ppt_slides(response_to_parse):
    # pexel = Pexels(api)
    prs = Presentation()
    slide_list = parse_response_to_ppt_content(response_to_parse)
    for i in range(len(slide_list)):
        # ppt_name = slide_list[i]['PPT File Name']
        title = slide_list[i]['Title']
        details = slide_list[i]['Details']
        st.header(title)
        st.write(details)

        layout = prs.slide_layouts[1]
        first_slide = prs.slides.add_slide(layout)  # Adding first slide
        first_slide.shapes.title.text = title
        first_slide.placeholders[1].text = details
        # if ppt_name != "":
        #     file_name = ppt_name + ".pptx"
        #     print(file_name)
        prs.save("generated.pptx")
        if i == (len(slide_list) - 1):
            st.title("AI_Generated_PPT.pptx created and saved to folder")
        # image = slide_list[i]['Image']
        # st.write(image)
        # search_photos = pexel.search_photos(query=image, orientation='', size='50x50', color='', locale='', page=1,
        #                                     per_page=1)
        # st.image(search_photos['photos'][0]['src']['original'])


st.title("🤖💻📑👔 AI Generated PPTs")
client = OpenAI()
prompt = ""

with st.form("ppt_generation_prompt"):
    topic = st.text_area("Enter topic related to which you want to generate slide upon")
    number_of_slides = st.slider("Enter number of slides")
    submitted = st.form_submit_button("Generate PPT")

    if submitted:
        prompt = (
                """Generate %s slide of content for PPT on %s. Follow below rules:
                    RULE-1: Generate content for the topic provided instead of suggesting how to create a PPT for that topic
                    RULE-2: Each slide must contain "Title:","Image: ","Details: "
                    RULE-3: For details in each slides keep content with bulletin points  
                    RULE-4: Each slide must be separated by "_________________________________"
                    RULE-5: Don't add slide number 
                    RULE-6: Differentiate content between 1000 and 5000 characters only HARD RULE : Follow all above rules""" % (
            number_of_slides, topic))
        print(prompt)

with st.spinner("Generating content..."):
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "system", "content": "You are a helpful assistant."},
                  {"role": "user", "content": prompt}]
        # prompt_text
    ).choices[0].message.content
    print(response)
    response_parser_response_to_ppt_slides(response)


